#!/usr/bin/env python3
"""Editable figures deck with the two schematic diagrams (framework, CVAE) as NATIVE, fully
editable PowerPoint shapes (rounded rectangles, cylinders, connectors with arrowheads, text),
followed by the data/map figures as images. -> paper/UrbanEV_figures_editable.pptx"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from PIL import Image

REPO = "/home/tomal/Documents/UrbanEV_Final_TRB/UrbanEV_Final_TRB"
FIG = os.path.join(REPO, "paper/figures")
OUT = os.path.join(REPO, "paper/UrbanEV_figures_editable.pptx")
SW, SH = Inches(13.333), Inches(7.5)

def C(h): return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
DATA, DATAF = C("5B6470"), C("EDEFF2")
GEN, GENF, GENB = C("2C6DA3"), C("DCE8F4"), C("EAF2FA")
ASN, ASNF, ASNB = C("3B8C5A"), C("DCEEE2"), C("EBF6EF")
SIM, SIMF, SIMB = C("2A8C93"), C("D6EEF0"), C("E8F5F6")
POL, POLF, POLB = C("C0872B"), C("FAEEDB"), C("FCF5E8")
INK, WHITE, GREY = C("1A1A1A"), C("FFFFFF"), C("555555")

prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def box(slide, x, y, w, h, title, sub=None, fill=WHITE, line=GEN, tcolor=None, tsize=12, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = line; sp.line.width = Pt(1.4); sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True; tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title; r.font.size = Pt(tsize); r.font.bold = True
    r.font.name = "Calibri"; r.font.color.rgb = tcolor or line
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub; r2.font.size = Pt(tsize - 3); r2.font.name = "Calibri"
        r2.font.color.rgb = INK
    return sp


def band(slide, x, y, w, h, fill, line, label):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill; sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    sp._element.addprevious(sp._element)  # keep; will be behind because added first
    # vertical label box on the far left
    lb = slide.shapes.add_textbox(Inches(x - 0.02), Inches(y), Inches(0.35), Inches(h))
    lb.rotation = 270
    tf = lb.text_frame; tf.word_wrap = False; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = line
    return sp


def arrow(slide, x1, y1, x2, y2, color=INK, width=1.4, elbow=False, dash=False):
    kind = MSO_CONNECTOR.ELBOW if elbow else MSO_CONNECTOR.STRAIGHT
    cn = slide.shapes.add_connector(kind, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color; cn.line.width = Pt(width)
    ln = cn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    if dash:
        d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'}); ln.append(d)
    return cn


def label(slide, x, y, w, text, size=9, color=GREY, italic=True, bold=False, align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.3))
    tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.italic = italic
    r.font.bold = bold; r.font.color.rgb = color; r.font.name = "Calibri"
    return tb


# =====================================================================
# SLIDE 1 — FRAMEWORK (native editable shapes)
# =====================================================================
s = prs.slides.add_slide(BLANK)
label(s, 0.3, 0.05, 12.7, "Figure 1. End-to-end modeling framework (editable)", size=16, color=C("0B2B4E"), italic=False, bold=True, align=PP_ALIGN.LEFT)
# bands (add first -> behind)
band(s, 0.55, 0.95, 12.5, 1.55, GENB, GEN, "GENERATIVE")
band(s, 0.55, 2.72, 12.5, 1.35, ASNB, ASN, "ASSIGNMENT")
band(s, 0.55, 4.25, 12.5, 1.05, SIMB, SIM, "SIMULATION")
band(s, 0.55, 5.45, 12.5, 1.35, POLB, POL, "POLICY")
# input cylinders
ins = [("RTS / MTS Survey", 1.0), ("MVA-2026 Reg.", 3.5), ("AFDC Chargers", 6.0), ("OSM POIs / Network", 8.5), ("NREL-278 Access", 11.0)]
for t, x in ins:
    box(s, x, 0.32, 2.1, 0.55, t, fill=DATAF, line=DATA, tcolor=C("333333"), tsize=8.5, shape=MSO_SHAPE.CAN)
    arrow(s, x + 1.05, 0.87, x + 1.05, 0.95, color=DATA, width=0.9, dash=True)
# generative boxes
box(s, 0.95, 1.25, 2.5, 1.0, "Household CVAE", "enc → z → dec", GENF, GEN, tsize=11)
box(s, 3.75, 1.25, 2.5, 1.0, "Trip CVAE (cond.)", "enc → z → dec", GENF, GEN, tsize=11)
box(s, 8.6, 1.25, 3.9, 1.0, "Synthetic Maryland", "population + activity chains", GENF, GEN, tsize=11)
arrow(s, 3.45, 1.75, 3.75, 1.75, GEN)
arrow(s, 6.25, 1.75, 8.6, 1.75, GEN)
# assignment boxes
axs = [(0.75, "EV Ownership Logit", "MVA-calibrated"), (3.75, "Vehicle & Charger", "make/model, NREL"),
       (6.75, "Geolocation", "OSM-POI"), (9.75, "MATSim Plans", "+ electric vehicles")]
for x, t, sub in axs:
    box(s, x, 3.0, 2.75, 0.85, t, sub, ASNF, ASN, tsize=10)
arrow(s, 10.55, 1.75, 10.55, 3.0, INK)                    # synth -> plans/assignment (down)
for x0 in [3.5, 6.5, 9.5]:
    arrow(s, x0, 3.42, x0 + 0.25, 3.42, ASN)
# simulation boxes
box(s, 0.95, 4.45, 6.3, 0.7, "UrbanEV / MATSim charging simulation", "range anxiety · walk · per-kWh cost", SIMF, SIM, tsize=10)
box(s, 8.6, 4.45, 3.9, 0.7, "Charging sessions", "+ base mobility (VMT)", SIMF, SIM, tsize=10)
arrow(s, 1.9, 3.85, 1.9, 4.45, INK)                        # plans... actually ownership col; keep vertical from assignment down
arrow(s, 8.6, 4.8, 7.25, 4.8, INK)                         # sessions -> UrbanEV
# policy boxes
pxs = [(0.95, "Shadow Gas-Tax Gap  R*", "iter-0 VMT × mpg⁻¹ × τ"),
       (5.05, "Recovery Instruments", "charging · registration · RUC"),
       (9.15, "Incidence & Equity", "Suits index · winners/losers")]
for x, t, sub in pxs:
    box(s, x, 5.7, 3.6, 0.85, t, sub, POLF, POL, tcolor=C("9A6A10"), tsize=11)
arrow(s, 2.1, 5.15, 2.1, 5.7, INK)                         # sim -> R*
arrow(s, 4.55, 6.12, 5.05, 6.12, POL)
arrow(s, 8.65, 6.12, 9.15, 6.12, POL)

# =====================================================================
# SLIDE 2 — CVAE ARCHITECTURE (native editable shapes)
# =====================================================================
s2 = prs.slides.add_slide(BLANK)
label(s2, 0.3, 0.05, 12.7, "Figure 2. Plain conditional VAE architecture (editable)", size=16, color=C("0B2B4E"), italic=False, bold=True, align=PP_ALIGN.LEFT)
LAT, LATF = C("C0872B"), C("FAEEDB")
HEAD, HEADF = C("7A4FA3"), C("ECE2F3")
COND, CONDF = C("B0466B"), C("F6E4EA")
# inputs
box(s2, 0.4, 1.3, 2.1, 1.0, "Categorical fields", "age, income, dwelling, …", DATAF, DATA, tcolor=C("333333"), tsize=10)
box(s2, 0.4, 2.6, 2.1, 0.9, "Numeric fields", "hh size, #vehicles", DATAF, DATA, tcolor=C("333333"), tsize=10)
box(s2, 2.8, 1.5, 1.2, 0.8, "embed", None, DATAF, DATA, tcolor=C("333333"), tsize=9)
box(s2, 4.2, 2.0, 0.55, 0.55, "⊕", None, WHITE, DATA, tcolor=INK, tsize=13, shape=MSO_SHAPE.OVAL)
# encoder
box(s2, 5.1, 1.3, 1.7, 1.9, "Encoder MLP", "Linear · LN+SiLU ×2", GENF, GEN, tsize=11)
# latent
box(s2, 7.0, 1.35, 1.0, 0.7, "μ", None, LATF, LAT, tsize=13)
box(s2, 7.0, 2.45, 1.0, 0.7, "log σ²", None, LATF, LAT, tsize=12)
box(s2, 8.25, 1.9, 0.6, 0.6, "~", None, WHITE, LAT, tsize=13, shape=MSO_SHAPE.OVAL)
box(s2, 9.1, 1.85, 0.95, 0.7, "z", None, LATF, LAT, tsize=13)
label(s2, 8.0, 2.65, 2.2, "z = μ + σ ⊙ ε,   ε ~ N(0, I)", size=8, color=GREY)
# decoder
box(s2, 10.3, 1.3, 1.7, 1.9, "Decoder MLP", "Linear · LN+SiLU · num branch", ASNF, ASN, tsize=11)
# heads + outputs
box(s2, 10.4, 3.5, 1.7, 0.7, "softmax head", "per categorical", HEADF, HEAD, tsize=9)
box(s2, 10.4, 4.35, 1.7, 0.7, "Gaussian head", "per numeric", ASNF, ASN, tsize=9)
box(s2, 12.15, 3.5, 0.9, 0.7, "x̂_cat", None, DATAF, DATA, tcolor=C("333333"), tsize=10)
box(s2, 12.15, 4.35, 0.9, 0.7, "x̂_num", None, DATAF, DATA, tcolor=C("333333"), tsize=10)
# condition + loss
box(s2, 4.7, 4.5, 3.0, 0.9, "condition c", "person + HH attrs (trip CVAE only)", CONDF, COND, tsize=10)
box(s2, 3.0, 5.7, 7.4, 0.9, "weighted ELBO", "Σ CE(x̂_cat) + Σ ½(x_num − x̂_num)²  +  β · KL( q(z|x,c) ‖ N(0,I) )", DATAF, DATA, tcolor=INK, tsize=11)
# arrows
arrow(s2, 2.5, 1.9, 2.8, 1.9, DATA)
arrow(s2, 4.0, 1.9, 4.2, 2.27, DATA)
arrow(s2, 2.5, 3.0, 4.2, 2.4, DATA)
arrow(s2, 4.75, 2.27, 5.1, 2.27, DATA)
arrow(s2, 6.8, 1.7, 7.0, 1.7, GEN)
arrow(s2, 6.8, 2.8, 7.0, 2.8, GEN)
arrow(s2, 8.0, 1.7, 8.25, 2.1, LAT)
arrow(s2, 8.0, 2.8, 8.25, 2.4, LAT)
arrow(s2, 8.85, 2.2, 9.1, 2.2, LAT)
arrow(s2, 10.05, 2.2, 10.3, 2.2, LAT)
arrow(s2, 11.15, 3.2, 11.15, 3.5, ASN)
arrow(s2, 11.25, 3.85, 11.25, 4.35, ASN)
arrow(s2, 12.1, 3.85, 12.15, 3.85, HEAD)
arrow(s2, 12.1, 4.7, 12.15, 4.7, ASN)
arrow(s2, 5.5, 4.5, 4.6, 2.55, COND, dash=True)            # cond -> encoder input
arrow(s2, 7.0, 4.5, 9.55, 2.55, COND, dash=True)           # cond -> z

# =====================================================================
# SLIDES 3+ — data/map figures (images with editable title/caption)
# =====================================================================
IMGS = [
    ("study_area.png", "Figure 3. Maryland study area", "EV owners by county (quantile choropleth) with 351 DC-fast chargers."),
    ("home_charger_access.png", "Figure 4. Home-charging access", "Assigned access by dwelling type and tenure (NREL-278 anchored); 91.8% overall."),
    ("val_tier1_population.png", "Figure 5. Tier 1 — population validation", "Marginal TVD, Cramér's V, an example joint, and the memorization (DCR) check."),
    ("val_tier2_trips.png", "Figure 6. Tier 2 — trip validation", "Distance, mode, departure hour, and daily VMT vs. survey."),
    ("val_tier3_ev.png", "Figure 7. Tier 3 — EV assignment validation", "County EV totals vs. MVA, income gradient, BEV share, make/model."),
    ("val_tier4_charging.png", "Figure 8. Tier 4 — charging validation", "Public diurnal profile vs. ChargePoint; venue mix by home access."),
    ("taxable_base.png", "Figure 9. Charging energy base by venue", "Only ~11% of energy is public (taxable)."),
    ("rate_revenue_frontier.png", "Figure 10. Surcharge revenue vs. rate", "Modeled surcharges fall far short of R*."),
    ("adequacy_equity_scatter.png", "Figure 11. Adequacy × equity of instruments", "No instrument is both fully adequate and progressive."),
]
DARK = C("0B2B4E")
for png, title, cap in IMGS:
    sl = prs.slides.add_slide(BLANK)
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.2), SW - Inches(1.0), Inches(0.7))
    p = tb.text_frame.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(22); r.font.bold = True; r.font.name = "Calibri"; r.font.color.rgb = DARK
    path = os.path.join(FIG, png)
    iw, ih = Image.open(path).size; ar = iw / ih
    w = Inches(11.2); h = Emu(int(w / ar))
    if h > Inches(5.3):
        h = Inches(5.3); w = Emu(int(h * ar))
    sl.shapes.add_picture(path, int((SW - w) / 2), Inches(1.05), width=int(w), height=int(h))
    cb = sl.shapes.add_textbox(Inches(0.6), SH - Inches(0.8), SW - Inches(1.2), Inches(0.6))
    cp = cb.text_frame; cp.word_wrap = True; pp = cp.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    rr = pp.add_run(); rr.text = cap; rr.font.size = Pt(12); rr.font.name = "Calibri"; rr.font.color.rgb = GREY

prs.save(OUT)
print(f"saved {OUT} | {len(prs.slides._sldIdLst)} slides (2 native-shape diagrams + 9 image figures)")
