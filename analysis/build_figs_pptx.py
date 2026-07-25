#!/usr/bin/env python3
"""Editable figures deck: one paper figure per slide, each with an editable title + caption
text box and the figure placed as a movable/resizable picture. -> paper/UrbanEV_figures_editable.pptx"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

REPO = "/home/tomal/Documents/UrbanEV_Final_TRB/UrbanEV_Final_TRB"
FIG = os.path.join(REPO, "paper/figures")
OUT = os.path.join(REPO, "paper/UrbanEV_figures_editable.pptx")
DARK = RGBColor(0x0B, 0x2B, 0x4E); GREY = RGBColor(0x55, 0x55, 0x55)
SW, SH = Inches(13.333), Inches(7.5)

FIGS = [
    ("framework_tikz-1.png", "Figure 1. End-to-end modeling framework",
     "Two plain CVAEs (generative) → assignment → UrbanEV/MATSim simulation → policy. Validation badges V1–V4."),
    ("cvae_tikz-1.png", "Figure 2. Plain conditional VAE architecture",
     "Embedded categorical + numeric inputs → MLP encoder → Gaussian latent (reparameterization) → decoder with softmax/Gaussian heads; weighted ELBO."),
    ("study_area.png", "Figure 3. Maryland study area",
     "EV owners by county (quantile choropleth) with 351 DC-fast chargers; 1,391 Level-2 omitted for legibility."),
    ("home_charger_access.png", "Figure 4. Home-charging access assignment",
     "Assigned home-charger access by dwelling type and tenure (NREL-278 anchored); overall 91.8%."),
    ("val_tier1_population.png", "Figure 5. Tier 1 — population synthesis validation",
     "Marginal TVD, Cramér's V associations, an example joint, and the memorization (DCR) check on held-out households."),
    ("val_tier2_trips.png", "Figure 6. Tier 2 — activity–travel validation",
     "Trip distance, travel mode, departure hour, and daily VMT vs. the weighted survey."),
    ("val_tier3_ev.png", "Figure 7. Tier 3 — EV ownership/assignment validation",
     "County EV totals vs. MVA-2026, the income–adoption gradient, BEV share, and fleet make/model."),
    ("val_tier4_charging.png", "Figure 8. Tier 4 — charging-behavior validation",
     "Simulated vs. observed public diurnal profile, and charging venue mix by home-charger access."),
    ("taxable_base.png", "Figure 9. Charging energy base by venue",
     "Only ~11% of the ~189 GWh annual charging energy is public (taxable); most is home charging on residential meters."),
    ("rate_revenue_frontier.png", "Figure 10. Surcharge revenue vs. rate",
     "Charging-surcharge revenue is the taxable base × rate; modeled surcharges fall far short of R* and the residual gap."),
    ("adequacy_equity_scatter.png", "Figure 11. Adequacy × equity of recovery instruments",
     "No instrument is both fully adequate and progressive; interstate-only RUC is the least-regressive adequate option."),
    ("recovery_waterfall.png", "Figure 12. Recovery waterfall",
     "R* = $33.3M → minus the existing Maryland fee → residual gap, with charging surcharges overlaid."),
    ("effective_tax_rate.png", "Figure 13. Effective tax rate by income",
     "Flat and registration fees are regressive: burden as a share of income falls with income."),
    ("suits_comparison.png", "Figure 14. Progressivity (Suits index) by instrument",
     "Instruments ranked by Suits index; charging surcharges and the flat fee are the most regressive."),
]

prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
blank = prs.slide_layouts[6]
missing = []
for png, title, cap in FIGS:
    s = prs.slides.add_slide(blank)
    # title (editable)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.2), SW - Inches(1.0), Inches(0.7))
    p = tb.text_frame.paragraphs[0]; p.text = title
    p.font.size = Pt(24); p.font.bold = True; p.font.name = "Calibri"; p.font.color.rgb = DARK
    # image (editable object)
    path = os.path.join(FIG, png)
    if os.path.isfile(path):
        iw, ih = Image.open(path).size; ar = iw / ih
        maxw, maxh = Inches(11.5), Inches(5.4)
        w = maxw; h = w / ar
        if h > maxh:
            h = maxh; w = h * ar
        s.shapes.add_picture(path, (SW - w) / 2, Inches(1.05), width=w, height=h)
    else:
        missing.append(png)
    # caption (editable)
    cb = s.shapes.add_textbox(Inches(0.6), SH - Inches(0.85), SW - Inches(1.2), Inches(0.7))
    cp = cb.text_frame; cp.word_wrap = True; pp = cp.paragraphs[0]; pp.text = cap
    pp.font.size = Pt(12); pp.font.name = "Calibri"; pp.font.color.rgb = GREY; pp.alignment = PP_ALIGN.CENTER

prs.save(OUT)
print(f"saved {OUT} | {len(prs.slides._sldIdLst)} slides | missing: {missing or 'none'}")
